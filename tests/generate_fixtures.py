"""Generate test fixture files for markitdown API verification."""
import csv
import json
import os
import sys
import zipfile

# Ensure we can import from the project venv
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

FIXTURES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fixtures")


def gen_csv():
    path = os.path.join(FIXTURES_DIR, "sample.csv")
    with open(path, "w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["Name", "Age", "City"])
        w.writerow(["Alice", "30", "Taipei"])
        w.writerow(["Bob", "25", "Tokyo"])
    return path


def gen_json():
    path = os.path.join(FIXTURES_DIR, "sample.json")
    with open(path, "w", encoding="utf-8") as f:
        json.dump({"project": "MarkItDown EZ", "version": "1.0", "features": ["convert", "preview"]}, f, indent=2)
    return path


def gen_xml():
    path = os.path.join(FIXTURES_DIR, "sample.xml")
    with open(path, "w", encoding="utf-8") as f:
        f.write('<?xml version="1.0" encoding="UTF-8"?>\n')
        f.write("<root>\n  <item id=\"1\">Hello</item>\n  <item id=\"2\">World</item>\n</root>\n")
    return path


def gen_html():
    path = os.path.join(FIXTURES_DIR, "sample.html")
    with open(path, "w", encoding="utf-8") as f:
        f.write("<!DOCTYPE html>\n<html><head><title>Test</title></head>\n")
        f.write("<body><h1>Hello MarkItDown</h1><p>This is a <strong>test</strong> paragraph.</p></body></html>\n")
    return path


def gen_docx():
    try:
        from docx import Document
    except ImportError:
        return None
    path = os.path.join(FIXTURES_DIR, "sample.docx")
    doc = Document()
    doc.add_heading("Test Document", level=1)
    doc.add_paragraph("This is a test paragraph for MarkItDown EZ.")
    doc.save(path)
    return path


def gen_xlsx():
    try:
        from openpyxl import Workbook
    except ImportError:
        return None
    path = os.path.join(FIXTURES_DIR, "sample.xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Score"])
    ws.append(["Alice", 95])
    ws.append(["Bob", 87])
    wb.save(path)
    return path


def gen_pptx():
    try:
        from pptx import Presentation
    except ImportError:
        return None
    path = os.path.join(FIXTURES_DIR, "sample.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = "This is a test presentation."
    prs.save(path)
    return path


def gen_pdf():
    """Generate a minimal valid PDF file."""
    path = os.path.join(FIXTURES_DIR, "sample.pdf")
    # Minimal PDF with text content
    content = """%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj
4 0 obj
<< /Length 44 >>
stream
BT /F1 12 Tf 100 700 Td (Hello PDF) Tj ET
endstream
endobj
5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj
xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000360 00000 n
trailer
<< /Size 6 /Root 1 0 R >>
startxref
441
%%EOF"""
    with open(path, "w") as f:
        f.write(content)
    return path


def gen_epub():
    """Generate a minimal EPUB (which is a ZIP with specific structure)."""
    path = os.path.join(FIXTURES_DIR, "sample.epub")
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("mimetype", "application/epub+zip")
        zf.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0"?>\n'
            '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">\n'
            '  <rootfiles>\n'
            '    <rootfile full-path="content.opf" media-type="application/oebps-package+xml"/>\n'
            '  </rootfiles>\n'
            '</container>',
        )
        zf.writestr(
            "content.opf",
            '<?xml version="1.0"?>\n'
            '<package xmlns="http://www.idpf.org/2007/opf" version="2.0">\n'
            '  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">\n'
            '    <dc:title>Test EPUB</dc:title>\n'
            '  </metadata>\n'
            '  <manifest>\n'
            '    <item id="ch1" href="chapter1.html" media-type="application/xhtml+xml"/>\n'
            '  </manifest>\n'
            '  <spine><itemref idref="ch1"/></spine>\n'
            '</package>',
        )
        zf.writestr(
            "chapter1.html",
            "<html><body><h1>Chapter 1</h1><p>Hello EPUB world.</p></body></html>",
        )
    return path


def gen_zip():
    """Generate a ZIP containing a text file and a CSV."""
    path = os.path.join(FIXTURES_DIR, "sample.zip")
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("readme.txt", "This is a readme inside a ZIP.")
        zf.writestr("data.csv", "col1,col2\nval1,val2\n")
    return path


if __name__ == "__main__":
    generators = {
        "CSV": gen_csv,
        "JSON": gen_json,
        "XML": gen_xml,
        "HTML": gen_html,
        "DOCX": gen_docx,
        "XLSX": gen_xlsx,
        "PPTX": gen_pptx,
        "PDF": gen_pdf,
        "EPUB": gen_epub,
        "ZIP": gen_zip,
    }
    for name, gen in generators.items():
        result = gen()
        if result:
            print(f"✅ {name}: {result}")
        else:
            print(f"❌ {name}: generator failed (missing dependency?)")
